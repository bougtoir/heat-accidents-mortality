#!/usr/bin/env python3
"""
Build the English manuscript (DOCX with inline figures + tables), an editable
figures PPTX (one figure per slide) and an editable tables DOCX.

Every numeric value is read from the generated result CSVs in data/processed/
and output/ -- nothing is hard-coded. Run `make all` first, then `make aap`.
"""
import os
import csv
import re
import shutil
import numpy as np
from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import parse_xml
from pptx import Presentation
from pptx.util import Inches as PInches, Pt as PPt

HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.dirname(HERE)
PROC = os.path.join(ROOT, "data", "processed")
OUT = os.path.join(ROOT, "output")
FIG = os.environ.get("FIGURES_DIR", os.path.join(OUT, "figures"))
MAN = os.environ.get("MANUSCRIPT_DIR", os.path.join(OUT, "manuscript"))
os.makedirs(MAN, exist_ok=True)


def read_csv(name):
    with open(os.path.join(PROC, name), encoding="utf-8") as f:
        return list(csv.DictReader(f))


def f(x, n=3):
    return f"{float(x):.{n}f}"


# ------------------------------------------------------------------ load
US = read_csv("us_attributable.csv")[0]
JP = read_csv("jp_attributable.csv")[0]
US_LAG = {r["window"]: r for r in read_csv("us_lag_response.csv")}
JP_LAG = {r["window"]: r for r in read_csv("jp_lag_response.csv")}
CDC = read_csv("cdc_heat_deaths.csv")
CTRL_ROWS = {r["model"]: r for r in read_csv("us_sensitivity_controls.csv")}
CTRL = (CTRL_ROWS.get("with_population_stateVMT_prcp_VIFscreened")
        or CTRL_ROWS.get("with_population_stateVMT_prcp")
        or CTRL_ROWS.get("with_national_VMT_gasoline")
        or list(CTRL_ROWS.values())[0])
CTRL_VMT = CTRL_ROWS.get("with_national_VMT_gasoline", CTRL)


def _ctrl_dropped_names():
    d = CTRL.get("dropped", "")
    if not d or d == "none" or str(d).lower() in ("nan", "none"):
        return "none"
    return ", ".join([x.split("(", 1)[0].strip() for x in str(d).split(";")])

CTRL_DROPPED_NAMES = _ctrl_dropped_names()
CTRL_KEPT_NAMES = "log-population offset and daily precipitation" if CTRL_DROPPED_NAMES != "none" else "all controls"
SUB = read_csv("us_subgroup_response.csv")
TOD = read_csv("us_timeofday_response.csv")
PROJ = read_csv("us_projection.csv")
USER = {r["group"]: r for r in SUB if r["dimension"] == "user"}
AGE = {r["group"]: r for r in SUB if r["dimension"] == "age"}
PROJ_D = {int(float(r["delta_degC"])): r for r in PROJ}
TOD_D = {r["hour_band"]: r for r in TOD}
TOD_MAX = max(TOD, key=lambda r: float(r["sameday_RR_+9C"]))
US_JP_RATIO = (float(US["total_deaths"]) / float(US["years"])) / (float(JP["total_deaths"]) / float(JP["years"]))
cdc_recent = [int(r["heat_deaths_X30"]) for r in CDC if 2016 <= int(r["year"]) <= 2020]
CDC_MEAN = np.mean(cdc_recent)
CDC_YRS = f"{min(int(r['year']) for r in CDC if 2016 <= int(r['year']) <= 2020)}-" \
          f"{max(int(r['year']) for r in CDC if 2016 <= int(r['year']) <= 2020)}"


def rr(d, key):
    return f"{f(d[key])} (95% CI {f(d[key+'_lo'])}-{f(d[key+'_hi'])})"


def lagrr(d):
    return f"{f(d['rr'])} ({f(d['lo'])}-{f(d['hi'])})"


# ------------------------------------------------------------------ references
# Vancouver: numbered automatically in order of first appearance in the text.
REF_TEXT = {
    "dlnm": "Gasparrini A, Armstrong B, Kenward MG. Distributed lag non-linear models. Stat Med 2010;29:2224-34.",
    "lancet": "Gasparrini A, Guo Y, Hashizume M, et al. Mortality risk attributable to high and low ambient temperature: a multicountry observational study. Lancet 2015;386:369-75.",
    "basu": "Basu R. High ambient temperature and mortality: a review of epidemiologic studies from 2001 to 2008. Environ Health 2009;8:40.",
    "fars": "National Highway Traffic Safety Administration. Fatality Analysis Reporting System (FARS). US Department of Transportation. https://www.nhtsa.gov/research-data/fatality-analysis-reporting-system-fars",
    "ghcn": "Menne MJ, Durre I, Vose RS, Gleason BE, Houston TG. An overview of the Global Historical Climatology Network-Daily database. J Atmos Ocean Technol 2012;29:897-910.",
    "cdc": "Centers for Disease Control and Prevention. CDC WONDER Underlying Cause of Death database. https://wonder.cdc.gov/",
    "npa": "National Police Agency of Japan. Traffic accident statistics open data. https://www.npa.go.jp/publications/statistics/koutsuu/opendata/",
    "attrib": "Gasparrini A, Leone M. Attributable risk from distributed lag models. BMC Med Res Methodol 2014;14:55.",
    "eia": "US Energy Information Administration. Weekly finished motor gasoline product supplied (PET.WGFUPUS2.W). https://www.eia.gov/",
    "fhwa": "Federal Highway Administration. Traffic Volume Trends. https://www.fhwa.dot.gov/policyinformation/travel_monitoring/tvt.cfm",
    "fhwa_vm2": "Federal Highway Administration. Highway Statistics VM-2: Annual state vehicle-miles travelled. https://www.fhwa.dot.gov/policyinformation/statistics.cfm",
    "census_pep": "US Census Bureau. Population Estimates Program. Annual state population estimates. https://www.census.gov/programs-surveys/popest.html",
    "ipcc": "IPCC. Climate Change 2021: The Physical Science Basis. Contribution of Working Group I to the Sixth Assessment Report of the Intergovernmental Panel on Climate Change. Cambridge: Cambridge University Press; 2021.",
    "daanen": "Daanen HAM, van de Vliert E, Huang X. Driving performance in cold, warm, and thermoneutral environments. Appl Ergon 2003;34:597-602.",
    "liang2022": "Liang M, Min M, Guo X, Song Q, Wang H, Li N, et al. The relationship between ambient temperatures and road traffic injuries: a systematic review and meta-analysis. Environ Sci Pollut Res 2022;29(33):50647-60.",
    "liang2021_aap": "Liang M, Zhao D, Wu Y, Ye P, Wang Y, Yao Z, et al. Short-term effects of ambient temperature and road traffic accident injuries in Dalian, Northern China: A distributed lag non-linear analysis. Accid Anal Prev 2021;153:106057.",
}
_ref_order = []


def cite(*keys):
    nums = []
    for k in keys:
        if k not in _ref_order:
            _ref_order.append(k)
        nums.append(str(_ref_order.index(k) + 1))
    # {..} markers are rendered as Word font superscript by _add_runs (Vancouver:
    # superscript numerals placed after the punctuation mark).
    return "{" + ",".join(nums) + "}"


# ------------------------------------------------------------------ docx helpers
def setup(doc):
    st = doc.styles["Normal"].font
    st.name = "Times New Roman"; st.size = Pt(11)


def h(doc, text, level=1):
    p = doc.add_heading(text, level=level)
    return p


def _add_runs(p, text):
    """Split text on {n} citation markers and render those as Word font
    superscript (editable in Word, unlike Unicode superscript glyphs)."""
    for seg in re.split(r"(\{[^}]+\})", text):
        if not seg:
            continue
        if seg[0] == "{" and seg[-1] == "}":
            r = p.add_run(seg[1:-1]); r.font.superscript = True
        else:
            p.add_run(seg)


def para(doc, text, space_after=6):
    p = doc.add_paragraph()
    _add_runs(p, text)
    p.paragraph_format.space_after = Pt(space_after)
    p.paragraph_format.line_spacing = 2.0
    return p


def labelled(doc, label, text, space_after=6, line_spacing=2.0):
    """Paragraph beginning with a bold run-in label (Lancet semi-structured
    abstract and Research-in-context panel)."""
    p = doc.add_paragraph()
    r = p.add_run(label); r.bold = True
    _add_runs(p, text)
    p.paragraph_format.space_after = Pt(space_after)
    p.paragraph_format.line_spacing = line_spacing
    return p


_EMBED = True  # when False, figures render as legend-only (separate-file submission)


def add_figure(doc, img, number, caption, width=5.6):
    if _EMBED:
        doc.add_paragraph()
        p = doc.add_paragraph(); p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.add_run().add_picture(os.path.join(FIG, img), width=Inches(width))
    cap = doc.add_paragraph()
    cap.paragraph_format.space_before = Pt(14)
    r = cap.add_run(f"Figure {number}. "); r.bold = True
    _add_runs(cap, caption)
    doc.add_paragraph()


def add_table(doc, number, title, header, rows):
    cap = doc.add_paragraph(); cap.paragraph_format.space_before = Pt(14)
    r = cap.add_run(f"Table {number}. "); r.bold = True
    _add_runs(cap, title)
    t = doc.add_table(rows=1, cols=len(header)); t.style = "Table Grid"
    for j, hd in enumerate(header):
        c = t.rows[0].cells[j].paragraphs[0].add_run(hd); c.bold = True
    for row in rows:
        cells = t.add_row().cells
        for j, v in enumerate(row):
            cells[j].text = str(v)
    doc.add_paragraph()


# ------------------------------------------------------------------ Word equations
# Native Office Math (OMML) so the equations are editable in Word, not LaTeX/images.
_M = "http://schemas.openxmlformats.org/officeDocument/2006/math"


def _mr(t):
    return f'<m:r><m:t xml:space="preserve">{t}</m:t></m:r>'


def _msub(base, sub):
    return f"<m:sSub><m:e>{base}</m:e><m:sub>{sub}</m:sub></m:sSub>"


def _mbar(inner):
    return ('<m:bar><m:barPr><m:pos m:val="top"/></m:barPr>'
            f"<m:e>{inner}</m:e></m:bar>")


def _mnary(sub, body):
    return ('<m:nary><m:naryPr><m:chr m:val="\u2211"/><m:limLoc m:val="subSup"/>'
            '<m:supHide m:val="1"/></m:naryPr>'
            f"<m:sub>{sub}</m:sub><m:sup/><m:e>{body}</m:e></m:nary>")


def add_equation(doc, inner):
    """Insert a centred, Word-native (OMML) display equation."""
    p = doc.add_paragraph(); p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    xml = (f'<m:oMathPara xmlns:m="{_M}"><m:oMath>{inner}</m:oMath></m:oMathPara>')
    p._p.append(parse_xml(xml))
    return p


# figures/tables metadata (used by both manuscript and pptx/tables docx)
FIGURES = [
    ("fig1_absolute_exposure_response.png",
     "United States absolute-temperature exposure-response for daily traffic-crash "
     "deaths (cumulative over lag 0-10 days), relative to the minimum-mortality "
     "temperature. The curve peaks at mild temperatures and declines at extreme heat, "
     "reflecting the seasonal driving-exposure envelope rather than an acute heat effect."),
    ("fig2_anomaly_exposure_response.png",
     "United States temperature-anomaly exposure-response (primary analysis): "
     "cumulative rate ratio of crash deaths versus the local day-of-year seasonal norm. "
     "Days hotter than normal carry higher crash mortality."),
    ("fig3_lag_response.png",
     "United States lag structure of the crash-death response to a +9 °C anomaly: an "
     "acute same-day excess followed by a 1-3 day deficit consistent with short-term "
     "displacement (harvesting) rather than a net addition of deaths."),
    ("fig_us_roaduser.png",
     "United States same-day rate ratio of crash death for a +9 °C anomaly by road-user "
     "type. Open-air users (motorcyclists, pedestrians, cyclists) show much larger heat "
     "effects than enclosed, often air-conditioned, vehicle occupants\u2014a gradient that "
     "tracks direct heat exposure rather than driving volume. "
     "This subgroup analysis is exploratory and is not adjusted for multiple comparisons."),
    ("fig_us_timeofday.png",
     "United States same-day rate ratio of crash death for a +9 °C anomaly by crash hour of "
     "day. The excess is largest for crashes in the hottest part of the day (12-17 h) and "
     "weakest in the cool morning (06-11 h). "
     "This subgroup analysis is exploratory and is not adjusted for multiple comparisons."),
    ("fig4_attributable_by_year.png",
     "United States estimated net heat-attributable crash deaths per year, "
     "2016-2022."),
    ("fig5_hidden_vs_official.png",
     "United States comparison of estimated net heat-attributable crash deaths per year "
     "with officially recorded direct-heat deaths (ICD-10 X30) from CDC WONDER."),
    ("fig_us_projection.png",
     "United States projected additional traffic-crash deaths per year under uniform "
     "warming of the daily temperature distribution by +1, +2 and +3 °C, holding driving "
     "activity and baseline rates constant. Bars show the point estimate with 95% "
     "Monte Carlo confidence intervals; this is a scenario projection, not an observed "
     "quantity."),
    ("jp_fig2_anomaly_exposure_response.png",
     "Japan temperature-anomaly exposure-response (2019-2024). The estimate is "
     "imprecise: the confidence band is wide and includes the null throughout."),
    ("cross_fig_us_vs_japan_sameday.png",
     "Same-day rate ratio of traffic-crash deaths for a +9 °C anomaly, United States "
     "versus Japan. The US shows a precise acute effect; the Japanese estimate is "
     "underpowered. This comparison is exploratory and is not adjusted for multiple comparisons."),
]

_BW_NOTES = {
    "fig_us_roaduser.png": "Black-and-white version: categories are distinguished by the y-axis labels.",
    "fig_us_timeofday.png": "Black-and-white version: the 12-17 h point is filled and the others are open to highlight the peak band.",
    "fig5_hidden_vs_official.png": "Black-and-white version: the right-hand (estimate) bar is hatched to distinguish it from the left-hand (official) bar.",
    "cross_fig_us_vs_japan_sameday.png": "Black-and-white version: US (square) and Japan (circle) estimates use different marker shapes."
}
if os.environ.get("FIGURES_BW") == "1":
    FIGURES = [(fn, cap + " " + _BW_NOTES.get(fn, "Black-and-white version: line styles and marker shapes distinguish categories.")) for fn, cap in FIGURES]


def tbl1(doc):
    add_table(doc, 1, "Panel description and model summary.",
              ["", "United States", "Japan"],
              [["Crash deaths", f"{int(float(US['total_deaths'])):,}", f"{int(float(JP['total_deaths'])):,}"],
               ["Years", int(float(US['years'])), int(float(JP['years']))],
               ["Spatial units", "50 states", "47 prefectures"],
               ["Dispersion (phi)", f(US['dispersion_phi'], 2), f(JP['dispersion_phi'], 2)]])


def tbl2(doc):
    add_table(doc, 2, "Rate ratio of traffic-crash deaths for a +9 °C temperature anomaly, by lag window.",
              ["Lag window (days)", "United States", "Japan"],
              [["0 (same day)", lagrr(US_LAG['lag0-0']), lagrr(JP_LAG['lag0-0'])],
               ["1-3", lagrr(US_LAG['lag1-3']), lagrr(JP_LAG['lag1-3'])],
               ["4-10", lagrr(US_LAG['lag4-10']), lagrr(JP_LAG['lag4-10'])],
               ["Cumulative 0-10", rr(US, 'cumRR_anom+9C'), rr(JP, 'cumRR_anom+9C')]])


def tbl3(doc):
    add_table(doc, 4, "United States net heat-attributable crash deaths versus officially recorded direct-heat deaths.",
              ["Quantity", "Estimate"],
              [["Net heat-attributable crash deaths / year", f"{float(US['net_heat_attributable_per_year']):.0f}"],
               ["Attributable fraction (%)", f(US['net_heat_attributable_fraction_pct'], 2)],
               [f"Official direct-heat deaths / year (X30, {CDC_YRS})", f"{CDC_MEAN:.0f}"]])


def subrr(d):
    return f"{f(d['sameday_RR_+9C'])} ({f(d['sameday_lo'])}-{f(d['sameday_hi'])})"


def ctrlrr(r):
    """Same-day RR from a sensitivity-controls row."""
    return f"{f(r['sameday_RR_anom+9C'])} ({f(r['sameday_RR_lo'])}-{f(r['sameday_RR_hi'])})"


def ctrlcum(r):
    """Cumulative RR from a sensitivity-controls row."""
    return f"{f(r['cumRR_anom+9C'])} ({f(r['cumRR_lo'])}-{f(r['cumRR_hi'])})"




def tbl4(doc):
    add_table(doc, 3, "United States same-day rate ratio of crash death for a +9 °C anomaly, "
              "by road-user type and age band (vulnerability analysis).",
              ["Subgroup", "Crash deaths", "Same-day RR (95% CI)"],
              [["Vehicle occupant", f"{int(USER['vehicle_occupant']['deaths']):,}", subrr(USER['vehicle_occupant'])],
               ["Motorcyclist", f"{int(USER['motorcyclist']['deaths']):,}", subrr(USER['motorcyclist'])],
               ["Pedestrian", f"{int(USER['pedestrian']['deaths']):,}", subrr(USER['pedestrian'])],
               ["Cyclist", f"{int(USER['cyclist']['deaths']):,}", subrr(USER['cyclist'])],
               ["Age <25 y", f"{int(AGE['<25']['deaths']):,}", subrr(AGE['<25'])],
               ["Age 25-64 y", f"{int(AGE['25-64']['deaths']):,}", subrr(AGE['25-64'])],
               ["Age 65+ y", f"{int(AGE['65+']['deaths']):,}", subrr(AGE['65+'])]])


def tbl5(doc):
    add_table(doc, 5, "United States projected additional traffic-crash deaths per year under "
              "uniform warming (activity held constant); scenario projection, not observed.",
              ["Warming", "Additional deaths/year (95% CI)"],
              [[f"+{d}\u00b0C",
                f"{float(PROJ_D[d]['extra_deaths_per_year']):.0f} "
                f"({float(PROJ_D[d]['extra_lo']):.0f}-{float(PROJ_D[d]['extra_hi']):.0f})"]
               for d in (1, 2, 3)])


SENS_LABELS = {
    "with_national_VMT_gasoline": "National VMT + motor gasoline",
    "with_population_offset": "State population offset",
    "with_population_stateVMT": "+ state annual VMT",
    "with_population_stateVMT_prcp": "+ state VMT + precipitation",
    "with_population_stateVMT_prcp_humidex": "+ humidex anomaly",
    "with_population_stateVMT_prcp_heat_index": "+ heat-index anomaly",
    "with_population_stateVMT_prcp_wbgt": "+ estimated WBGT anomaly",
    "with_population_stateVMT_prcp_VIFscreened": "VIF-screened full controls",
}
SENS_ORDER = [
    "with_national_VMT_gasoline",
    "with_population_offset",
    "with_population_stateVMT",
    "with_population_stateVMT_prcp",
    "with_population_stateVMT_prcp_VIFscreened",
    "with_population_stateVMT_prcp_humidex",
    "with_population_stateVMT_prcp_heat_index",
    "with_population_stateVMT_prcp_wbgt",
]


def _vif_cell(r):
    v = r.get("max_control_VIF")
    if v in ("", "nan", "NaN", None) or (isinstance(v, float) and np.isnan(v)):
        return "-"
    try:
        return f"{float(v):.1f}"
    except (TypeError, ValueError):
        return str(v)


def tbl6(doc):
    rows = []
    for key in SENS_ORDER:
        if key not in CTRL_ROWS:
            continue
        r = CTRL_ROWS[key]
        rows.append([
            SENS_LABELS.get(key, key),
            ctrlrr(r),
            ctrlcum(r),
            _vif_cell(r),
        ])
    add_table(doc, 6, "United States sensitivity of the +9 \u00b0C anomaly association to activity, "
              "precipitation and heat-stress controls. All added continuous controls were z-scored; "
              "population entered as a log offset; heat-stress metrics were anomaly-derived. The "
              "VIF-screened row is the primary full-controls model, selected by iteratively removing "
              "added controls with VIF > 5. Heat-stress metrics were tested individually because they "
              "are functions of temperature and can be collinear with the anomaly.",
              ["Model", "Same-day RR (95% CI)", "Cumulative RR (95% CI)", "Max control VIF"],
              rows)


def build_manuscript(filename="heat_crash_mortality.docx", embed=True):
    global _EMBED, _ref_order
    _EMBED = embed
    _ref_order = []  # reset Vancouver numbering for each build
    doc = Document(); setup(doc)

    t = doc.add_paragraph(); t.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = t.add_run("Ambient heat as an under-recognised risk factor for US traffic-crash "
                  "mortality: a distributed-lag analysis with road-safety implications")
    r.bold = True; r.font.size = Pt(14)

    ap = doc.add_paragraph(); ap.alignment = WD_ALIGN_PARAGRAPH.CENTER
    ap.add_run("Tatsuki Onishi")
    af = doc.add_paragraph(); af.alignment = WD_ALIGN_PARAGRAPH.CENTER
    af.add_run("Data Science and AI Innovation Research Promotion Center, "
               "Shiga University of Medical Science, Seta Tsukinowa-cho, "
               "Otsu, Shiga 520-2192, Japan")
    cp = doc.add_paragraph(); cp.alignment = WD_ALIGN_PARAGRAPH.CENTER
    cp.add_run("Correspondence to: Tatsuki Onishi, Data Science and AI Innovation "
               "Research Promotion Center, Shiga University of Medical Science, "
               "Seta Tsukinowa-cho, Otsu, Shiga 520-2192, Japan. "
               "E-mail: bougtoir@gmail.com. ORCID: [iD to be added]")

    h(doc, "Summary", 1)
    labelled(doc, "Background ",
             "Ambient heat is an established mortality risk factor; because driving "
             "continues in heat, fatal crashes may involve heat-related impairment or "
             "unrecorded heat illness. We quantified the acute association between local "
             "heat anomalies and daily traffic-crash mortality in the US and Japan.")
    labelled(doc, "Methods ",
             "We fitted quasi-Poisson distributed-lag models of local seasonal "
             "temperature anomalies to US state-day and Japan prefecture-day panels, "
             "adjusting for spatial, seasonal, trend and day-of-week effects. Data were FARS "
             "(2016-2022), Japanese NPA open data (2019-2024) and GHCN-Daily.")
    labelled(doc, "Findings ",
             f"In the USA ({int(float(US['total_deaths'])):,} crash deaths), a +9\u00b0C anomaly "
             f"raised same-day crash mortality, RR {rr(US, 'sameday_RR_anom+9C')}, "
             "and remained similar in a VIF-screened model adjusting for state population (offset) "
             f"and precipitation (VIF > 5 iteratively dropped {CTRL_DROPPED_NAMES}; RR "
             f"{ctrlrr(CTRL)}; cumulative RR {ctrlcum(CTRL)}). The excess was larger "
             f"for open-air users (motorcyclists RR {f(USER['motorcyclist']['sameday_RR_+9C'])}, "
             f"pedestrians RR {f(USER['pedestrian']['sameday_RR_+9C'])}) than vehicle occupants "
             f"(RR {f(USER['vehicle_occupant']['sameday_RR_+9C'])}). Net heat-attributable deaths "
             f"were {float(US['net_heat_attributable_per_year']):.0f} per year "
             f"({f(US['net_heat_attributable_fraction_pct'],2)}%), similar to "
             f"{CDC_MEAN:.0f} recorded direct-heat deaths (ICD-10 X30); "
             "+1 to +3\u00b0C warming projected "
             f"{float(PROJ_D[1]['extra_deaths_per_year']):.0f}-"
             f"{float(PROJ_D[3]['extra_deaths_per_year']):.0f} additional deaths per year. "
             f"Japan data ({int(float(JP['total_deaths'])):,} deaths) were too sparse for a "
             "precise estimate.")
    labelled(doc, "Interpretation ",
             "Unusually hot days are associated with an acute excess of US crash "
             "deaths comparable to recorded direct-heat mortality, concentrated in "
             "heat-exposed road users and projected to grow with warming. This is "
             "consistent with an under-recognised heat contribution, but cannot establish "
             "heat illness in any individual crash and open-air gradient may partly "
             "reflect weather-related activity. Findings suggest heat-aware road safety "
             "and climate adaptation planning could address an uncounted burden.")
    labelled(doc, "Funding ", "None.")

    h(doc, "Research in context", 1)
    labelled(doc, "Evidence before this study ",
             "Ambient heat is a well-established driver of daily mortality, and "
             "distributed-lag and distributed-lag non-linear models linking temperature to "
             "death are a mature methodology. Prior work has largely focused on "
             "cardiovascular, respiratory, and all-cause mortality; heat illness is thought "
             "to be substantially under-recorded because it is rarely investigated after "
             "death. Traffic-crash deaths, which seldom undergo a heat-oriented post-mortem "
             "examination, have received little attention as a possible reservoir of "
             "unrecognised heat mortality. We are not aware of national population-level "
             "studies quantifying the acute crash-mortality burden of hotter-than-normal "
             "days alongside officially recorded direct-heat deaths.")
    labelled(doc, "Added value of this study ",
             "Using only public data, we show at the population level that days hotter than "
             "the local seasonal norm carry an acute excess of US traffic-crash deaths that "
             "persists after adjustment for national driving activity, is concentrated in "
             "heat-exposed open-air road users, and is comparable in magnitude to all "
             "officially recorded direct-heat mortality. We provide scenario projections "
             "showing this hidden burden would grow under uniform warming, and we report a "
             "much weaker, imprecise signal in Japan.")
    labelled(doc, "Implications of all the available evidence ",
             "The results suggest that routine cause-of-death coding may miss a heat "
             "contribution to road deaths of societal importance, and that this burden is "
             "sensitive to future warming. Because the analysis is ecological, it cannot "
             "attribute any individual crash to heat illness and cannot exclude "
             "weather-related changes in activity; it motivates individual-level "
             "investigation, such as heat-aware forensic assessment of crash fatalities.")

    h(doc, "Introduction", 1)
    para(doc,
         "Road traffic crashes are a leading cause of preventable death, and climate change "
         "is adding a heat-related layer of risk. Heat waves are becoming more frequent, and "
         "ambient heat is an established driver of mortality"
         f".{cite('lancet','basu')} Recent syntheses also link higher ambient "
         f"temperature to increased road traffic accidents and injuries.{cite('liang2022','liang2021_aap')} Occupational and driving activity "
         "do not stop during heat, so some fatal accidents plausibly involve heat-related "
         "impairment or overt heat illness. Traffic-crash deaths are of particular interest "
         "because decedents rarely undergo the post-mortem assessment that would identify a "
         "heat contribution, so any such contribution would be largely invisible in routine "
         "cause-of-death coding. We therefore ask whether days that are hotter than the "
         "local seasonal norm carry excess traffic-crash mortality, and how the magnitude "
         "compares with officially recorded direct-heat deaths. This is an ecological, "
         "population-level study whose aim is to gauge the potential societal burden of "
         "under-recognised heat illness\u2014the deaths that would be missing from official "
         "heat statistics\u2014rather than to diagnose heat illness in any individual crash. "
         "If population-level heat contributes to road deaths, heat-aware road safety "
         "messaging and occupational protections for outdoor riders and drivers could "
         "reduce a burden currently invisible to both heat-mortality and road-safety "
         "surveillance.")

    h(doc, "Methods", 1)
    para(doc,
         "We built daily panels of traffic-crash deaths by US state (Fatality Analysis "
         f"Reporting System [FARS], 2016-2022){cite('fars')} and by Japanese prefecture "
         f"(National Police Agency [NPA] accident open data, 2019-2024).{cite('npa')} "
         "Daily mean temperature was taken from the Global Historical Climatology "
         f"Network-Daily (GHCN-Daily) stations{cite('ghcn')} as the mean of TMAX and TMIN, "
         "aggregated to each spatial unit from its nearest reporting stations. State-level "
         "precipitation (PRCP) and, in sensitivity analyses, daily average dew point (ADPT), "
         "relative humidity (RHAV) and average wet-bulb temperature (AWBT) were taken from the "
         "same GHCN-Daily stations and used to compute humidex, the US National Weather Service "
         "heat index (converted back to degrees Celsius) and an estimated wet-bulb globe "
         "temperature (0.7 times wet-bulb plus 0.3 times air temperature). For every unit "
         "we estimated a cyclic day-of-year climatology and defined the temperature anomaly "
         "as the observed temperature minus that climatology,", space_after=2)
    add_equation(doc,
                 _msub(_mr("A"), _mr("u,t")) + _mr(" = ") + _msub(_mr("T"), _mr("u,t"))
                 + _mr(" \u2212 ") + _msub(_mbar(_mr("T")), _mr("u")) + _mr("(")
                 + _msub(_mr("d"), _mr("t")) + _mr(")"))
    para(doc,
         "for the anomaly A of spatial unit u on day t, where T is the observed daily mean "
         "temperature and the barred T is the day-of-year climatology d(t) for that unit.",
         space_after=6)
    para(doc,
         "Absolute temperature is strongly confounded by the seasonal cycle of driving "
         "exposure, so our primary model used the anomaly. We fitted a quasi-Poisson "
         f"distributed-lag model{cite('dlnm')} of the form", space_after=2)
    add_equation(doc,
                 _mr("log E(") + _msub(_mr("Y"), _mr("u,t")) + _mr(") = ")
                 + _mnary(_mr("k"), _msub(_mr("f"), _mr("k")) + _mr("(")
                          + _msub(_mr("A"), _mr("u,t")) + _mr(")"))
                 + _mr(" + ") + _msub(_mr("\u03b1"), _mr("u"))
                 + _mr(" + ") + _msub(_mr("s"), _mr("r(u)")) + _mr("(")
                 + _msub(_mr("d"), _mr("t")) + _mr(")")
                 + _mr(" + h(t) + ") + _msub(_mr("\u03b4"), _mr("dow(t)")))
    para(doc,
         "where Y is the daily crash-death count, the f\u2096 are constant-free natural cubic "
         "splines of the anomaly within three lag windows (same day, 1-3 days, 4-10 days) "
         "summed over k, \u03b1\u1d64 are spatial fixed effects, s is a region-specific cyclic "
         "seasonal spline of day of year, h a long-term time trend and \u03b4 a day-of-week "
         "effect. Each exposure-response spline had 4 df (3 effective df after removing the "
         "constant); the US and Japanese seasonal splines had 8 and 6 df respectively, and "
         "the long-term trend had 3 df per study year. Overdispersion was handled with a "
         "quasi-Poisson dispersion parameter. Net heat-attributable deaths were computed by "
         f"the method of Gasparrini and Leone{cite('attrib')} with Monte Carlo (MC) confidence "
         "intervals (CIs). As sensitivity analyses for the US we added (i) national "
         f"vehicle-miles travelled (VMT){cite('fhwa')} and finished-motor-gasoline product supplied"
         f"{cite('eia')} as activity proxies, and (ii) a richer control set with state annual "
         f"population entered as a log offset{cite('census_pep')}, state annual VMT"
         f"{cite('fhwa_vm2')} and daily state precipitation, all as z-scored linear terms except "
         "the population offset, plus GHCN-derived heat-stress metrics (humidex, NOAA heat index, "
         "and an estimated wet-bulb globe temperature) expressed as anomalies relative to their "
         "day-of-year climatology. Because heat-stress metrics are functions of temperature and can "
         "be collinear with the anomaly exposure, we selected the primary full-controls model by a "
         "variance-inflation-factor (VIF) screen (iteratively dropping added controls with VIF > 5), "
         "and tested the heat-stress metrics individually as sensitivity analyses. "
         "State fixed effects capture time-invariant cross-state differences in baseline risk, while the "
         "log-population offset scales expected counts within each state across years; because state "
         "annual VMT is correlated with population size, the two partly adjust for the same scale factor, "
         "so their simultaneous use is reported as a sensitivity analysis rather than the primary specification. "
         "Officially recorded direct-heat deaths "
         "(International Classification of Diseases, 10th revision [ICD-10] code X30) were "
         f"obtained from CDC WONDER.{cite('cdc')} We performed three further US analyses: (i) "
         "refitting the same-day model within crash-hour bands to test whether the excess "
         "concentrates in the hottest hours; (ii) refitting it by road-user type (vehicle "
         "occupant, motorcyclist, pedestrian, cyclist) and age band from FARS person records, "
         "as a mechanistic and vulnerability analysis; and (iii) projecting the additional "
         "crash deaths under uniform warming scenarios of the daily temperature distribution "
         "(+1, +2, +3 °C), holding activity and baseline rates constant. All data are public "
         "and the full pipeline is reproducible from source (see Data availability). Reporting "
         "follows the Strengthening the Reporting of Observational Studies in Epidemiology "
         "(STROBE) guideline (checklist provided).")
    para(doc,
         "The +9 °C anomaly level was chosen because it lies near the upper end of the observed "
         "positive-anomaly range while remaining within the support of the estimated anomaly "
         "response (approximately -12 to +10 °C), so the estimate does not require extrapolation. "
         "The quasi-Poisson variance accounts for overdispersion but does not model residual "
         "temporal autocorrelation explicitly; the distributed-lag, seasonal, trend and "
         "day-of-week terms absorb most short-term serial correlation, but any remaining "
         "autocorrelation could affect confidence intervals (see Limitations).")
    para(doc,
         "No ethics approval was required because the study used publicly available, "
         "aggregated and de-identified data with no individual patient records.")
    para(doc,
         "Generative artificial-intelligence tools were used to assist with code development "
         "and manuscript drafting; the author reviewed and verified all content and takes full "
         "responsibility for the final work.")

    h(doc, "Results", 1)
    para(doc,
         "The panels comprised "
         f"{int(float(US['total_deaths'])):,} US crash deaths over {int(float(US['years']))} "
         f"years (50 states) and {int(float(JP['total_deaths'])):,} Japanese crash deaths over "
         f"{int(float(JP['years']))} years (47 prefectures); descriptive characteristics and "
         "the quasi-Poisson dispersion are summarised in Table 1.")
    tbl1(doc)
    para(doc,
         "In the descriptive absolute-temperature model, US crash mortality peaked at mild "
         "temperatures and declined at extreme heat (Fig. 1); this reflects the seasonal "
         "driving-exposure envelope and cannot be read as an acute heat effect. When we "
         "instead used the temperature anomaly, days hotter than the seasonal norm carried "
         f"higher crash mortality (Fig. 2). A +9 °C anomaly was associated with higher same-day crash deaths, "
         f"RR {rr(US, 'sameday_RR_anom+9C')}, with a cumulative lag 0-10 RR of "
         f"{rr(US, 'cumRR_anom+9C')}.")
    add_figure(doc, FIGURES[0][0], 1, FIGURES[0][1])
    add_figure(doc, FIGURES[1][0], 2, FIGURES[1][1])
    para(doc,
         "The lag structure showed an acute same-day excess followed by a deficit at 1-3 "
         f"days, RR {lagrr(US_LAG['lag1-3'])}, consistent with short-term displacement "
         "(harvesting) rather than a net addition of deaths (Fig. 3; Table 2). The "
         "same-day association was essentially unchanged after adjusting for national "
         f"vehicle-miles travelled and gasoline supplied (RR {f(CTRL_VMT['sameday_RR_anom+9C'])}, "
         f"95% CI {f(CTRL_VMT['sameday_RR_lo'])}-{f(CTRL_VMT['sameday_RR_hi'])}; "
         f"cumulative RR {ctrlcum(CTRL_VMT)}). A variance-inflation-factor (VIF) screened full-controls "
         "model—starting with state population as a log offset plus z-scored state annual VMT, daily "
         "precipitation and the three heat-stress metrics, then iteratively dropping any added control "
         f"with VIF > 5—retained only the {CTRL_KEPT_NAMES}. "
         f"Iteratively dropped variables were {CTRL_DROPPED_NAMES}. "
         "Its same-day RR remained similar "
         f"(RR {ctrlrr(CTRL)}, 95% CI {f(CTRL['sameday_RR_lo'])}-{f(CTRL['sameday_RR_hi'])}), "
         f"but the cumulative 0-10 day RR was attenuated (RR {ctrlcum(CTRL)}). "
         "Heat-stress metrics were tested individually because they are functions of temperature and can be "
         "collinear with the anomaly; they yielded variable same-day estimates and the heat-index model "
         "attenuated toward the null.")
    add_figure(doc, FIGURES[2][0], 3, FIGURES[2][1])
    tbl2(doc)
    para(doc,
         "Two exploratory analyses were consistent with a heat mechanism, although neither "
         "can exclude residual confounding. First, the same-day excess differed by "
         f"road-user exposure (Fig. 4; Table 3). The effect was small for enclosed, often "
         f"air-conditioned vehicle occupants, RR {subrr(USER['vehicle_occupant'])}, but much "
         f"larger for open-air users: motorcyclists, RR {subrr(USER['motorcyclist'])}; "
         f"pedestrians, RR {subrr(USER['pedestrian'])}; and cyclists, RR "
         f"{subrr(USER['cyclist'])}, a pattern consistent with direct bodily heat "
         f"exposure and physical exertion.{cite('daanen')} This gradient could also partly "
         "reflect greater discretionary open-air travel (motorcycling, cycling, walking) on "
         "hotter-than-normal days, which our aggregate activity proxies do not capture by "
         "mode. Across age bands the excess was similar: age <25 y RR "
         f"{subrr(AGE['<25'])}; 25-64 y RR {subrr(AGE['25-64'])}; 65+ y RR "
         f"{subrr(AGE['65+'])} (Table 3), giving no clear age gradient. Second, the excess "
         "was largest for crashes in the hottest part of the day "
         f"({TOD_MAX['hour_band']} h), RR {subrr(TOD_MAX)}, and weakest in the cool morning "
         f"(06-11 h), RR {subrr(TOD_D['06-11'])}; the overnight band was also elevated "
         f"(00-05 h), RR {subrr(TOD_D['00-05'])}, so the diurnal pattern is not a clean "
         "daytime-only gradient (Fig. 5). These subgroup analyses are hypothesis-generating: "
         "they involve multiple unadjusted comparisons and are not individual-level tests "
         "of susceptibility or mechanism.")
    add_figure(doc, FIGURES[3][0], 4, FIGURES[3][1])
    tbl4(doc)
    add_figure(doc, FIGURES[4][0], 5, FIGURES[4][1])
    para(doc,
         "Net heat-attributable crash deaths were "
         f"{float(US['net_heat_attributable_per_year']):.0f} per year (95% CI "
         f"{float(US['net_heat_attributable_lo'])/float(US['years']):.0f}-"
         f"{float(US['net_heat_attributable_hi'])/float(US['years']):.0f}), an attributable "
         f"fraction of {f(US['net_heat_attributable_fraction_pct'],2)}% (Fig. 6; Table 4). This is of "
         f"the same order as the {CDC_MEAN:.0f} officially recorded direct-heat deaths per "
         f"year (ICD-10 X30, {CDC_YRS}; Fig. 7), so a heat contribution to road deaths on the "
         "scale of recorded direct-heat mortality could go unrecorded in cause-of-death "
         "statistics.")
    add_figure(doc, FIGURES[5][0], 6, FIGURES[5][1])
    tbl3(doc)
    add_figure(doc, FIGURES[6][0], 7, FIGURES[6][1])
    para(doc,
         "To gauge the trajectory under continued warming, we applied the estimated anomaly "
         "exposure-response to uniformly warmer daily temperatures, holding driving activity "
         f"and baseline rates constant.{cite('ipcc')} A uniform +1 °C shift projected "
         f"{float(PROJ_D[1]['extra_deaths_per_year']):.0f} additional US crash deaths per "
         f"year (95% CI {float(PROJ_D[1]['extra_lo']):.0f}-{float(PROJ_D[1]['extra_hi']):.0f}), "
         f"rising to {float(PROJ_D[3]['extra_deaths_per_year']):.0f} "
         f"({float(PROJ_D[3]['extra_lo']):.0f}-{float(PROJ_D[3]['extra_hi']):.0f}) at +3 °C "
         "(Fig. 8; Table 5). These are scenario projections, not observed quantities, and "
         "assume the anomaly response is stable under a warmer mean climate.")
    add_figure(doc, FIGURES[7][0], 8, FIGURES[7][1])
    tbl5(doc)
    para(doc,
         "Sensitivity to activity, precipitation and heat-stress controls is summarized in Table 6. "
         "All added continuous terms were z-scored linear terms; state population entered as a log offset, "
         "and heat-stress metrics were expressed as anomalies relative to their day-of-year climatology. "
         "The primary full-controls model was selected by iteratively removing added controls with VIF > 5; "
         f"dropped variables were {CTRL.get('dropped', 'none')}. "
         "Heat-stress metrics showed high collinearity with the temperature anomaly when entered individually, "
         "so they are reported separately as sensitivity analyses rather than being included simultaneously.")
    tbl6(doc)
    h(doc, "Exploratory external validation: Japan comparison", 2)
    para(doc,
         "We did not pool the Japanese estimates with the US estimates; the comparison is "
         "presented only as an exploratory external validation. "
         f"In Japan ({int(float(JP['total_deaths'])):,} crash deaths over "
         f"{int(float(JP['years']))} years) the anomaly exposure-response was imprecise, with "
         "a wide confidence band that included the null throughout (Fig. 9); the same-day "
         f"point estimate for a +9 °C anomaly was {rr(JP, 'sameday_RR_anom+9C')} and was not "
         "stable across specifications. A direct comparison of the two countries (Fig. 10) shows a precise acute effect in the "
         f"US, whereas the Japanese panel—with roughly {US_JP_RATIO:.0f}-fold fewer annual crash deaths and sparser temperature coverage—is underpowered.")
    add_figure(doc, FIGURES[8][0], 9, FIGURES[8][1])
    add_figure(doc, FIGURES[9][0], 10, FIGURES[9][1])

    h(doc, "Discussion", 1)
    para(doc,
         f"A +9 °C temperature anomaly was associated with a same-day rate ratio of {rr(US, 'sameday_RR_anom+9C')} "
         f"for US traffic-crash deaths, and the net annual heat-attributable burden was "
         f"{float(US['net_heat_attributable_per_year']):.0f} deaths (95% CI "
         f"{float(US['net_heat_attributable_lo'])/float(US['years']):.0f}-"
         f"{float(US['net_heat_attributable_hi'])/float(US['years']):.0f}), "
         f"similar to the {CDC_MEAN:.0f} officially recorded direct-heat deaths per year. "
         "This acute, same-day excess was not explained by aggregate driving activity, and "
         "it remained in the VIF-screened full-controls model (added controls iteratively removed "
         f"when VIF > 5, retaining only the {CTRL_KEPT_NAMES}; dropped: {CTRL_DROPPED_NAMES}), "
         f"RR {ctrlrr(CTRL)}, 95% CI {f(CTRL['sameday_RR_lo'])}-{f(CTRL['sameday_RR_hi'])}). "
         f"The cumulative 0-10 day RR in that model was attenuated (RR {ctrlcum(CTRL)}). "
         "The 1-3 day deficit "
         "indicates that part of the excess reflects short-term displacement (harvesting). These "
         "findings are consistent with heat-related impairment or under-recognised heat illness "
         f"contributing to road deaths without appearing in cause-of-death data.{cite('liang2022','liang2021_aap')}")
    para(doc,
         "Two features of the data argue against confounding by overall driving volume and "
         "are consistent with a direct heat effect. The excess rose with direct heat exposure: "
         "it was small for enclosed, often air-conditioned vehicle occupants but several-fold "
         "larger for motorcyclists, pedestrians and cyclists, who are directly exposed and "
         "often physically exerting. It was also largest in the hottest hours of the day. Both "
         "patterns are consistent with heat degrading psychomotor and cognitive performance. "
         "These gradients are not proof of mechanism, however: open-air travel is itself "
         "weather-sensitive, so more motorcycling, cycling and walking on hotter-than-normal "
         "days could inflate the open-air estimates through greater exposure rather than "
         "physiology, and our activity proxies are national and not mode-specific. Applied to a "
         "warmer mean climate, the same exposure-response implies hundreds of additional US "
         "crash deaths per year per degree of warming; because these deaths would continue to "
         "be coded as ordinary crashes, the societal heat burden they represent would remain "
         "uncounted by heat-mortality surveillance.")
    h(doc, "Implications for road safety", 2)
    para(doc,
         "These results point to practical implications for road safety and climate "
         "adaptation. Because the same-day excess was largest among open-air road users "
         "and in the hottest hours, heat warnings and travel advice could be targeted at "
         "motorcyclists, cyclists, pedestrians and outdoor delivery riders during hot "
         "periods. Shared e-scooter and e-bike fleets and public bikeshare systems should also "
         "be considered: batteries, motors and tyres are heat-stressed, increasing the risk of "
         "sudden power loss or brake failure, while docking and charging points without shade "
         "expose users and maintenance staff to thermal load during check-out, parking and "
         "battery swap. Service operators could temporarily reduce or suspend rentals above "
         "high-temperature thresholds, relocate or shade docking and charging infrastructure, "
         "and integrate heat warnings into app-based routing. Linking shared-mobility heat "
         "alerts with public transit and cooled shelter maps would help preserve access for "
         "users who depend on these modes, particularly in low-income neighbourhoods where "
         "private air-conditioned transport is less available. Road safety agencies typically "
         "classify fatal crashes by driver error, vehicle failure or road conditions; the "
         "possibility that ambient heat impairs psychomotor or cognitive performance is rarely "
         "considered. Integrating temperature forecasts into crash-prevention messaging and "
         "ensuring heat-health plans reach outdoor workers and commuters is consistent with a "
         "Safe System approach to environmental risk and could reduce a burden that is currently "
         "invisible to both heat-mortality and road-safety surveillance, though intervention "
         "studies are needed to test whether such heat-aware measures reduce crash risk.")
    para(doc,
         "Several limitations apply. First, this is an ecological, population-level "
         "association: it cannot establish that heat caused illness in any specific crash, "
         "and neither FARS nor NPA records post-mortem heat diagnosis. Second, exposure is "
         "measured at the state/prefecture level from the nearest GHCN-Daily station; this "
         "misclassifies individual exposure and does not capture sub-unit variation such as "
         "urban heat islands. We therefore added state population as an offset and evaluated "
         "GHCN-derived humidity, heat-index and estimated wet-bulb globe temperature metrics in "
         "sensitivity models, but these station-based heat-stress measures still do not measure "
         "personal exposure. Third, activity adjustment used national monthly proxies for the main "
         "model and annual state VMT for the sensitivity model; neither is daily, mode-specific or "
         "shared-mobility-specific, and daily state-level or mode-specific activity could alter "
         "the road-user and time-of-day gradients. In particular, publicly available daily "
         "national counts of pedestrian, bicycle and shared-mobility trips are not available, so "
         "we cannot directly separate a weather-related activity shift from a physiological or "
         "mechanical heat effect for those modes. Fourth, the Japanese analysis is an "
         "exploratory external validation only; its shorter series, far smaller death counts "
         "and sparser GHCN-Daily coverage preclude a precise estimate, and we do not pool "
         "the two countries. Fifth, the warming projection is a scenario calculation that "
         "holds driving activity, the vehicle fleet and behaviour fixed and assumes the "
         "anomaly response is stable under a warmer mean climate; it should be read as "
         "illustrative of magnitude, not as a forecast. Sixth, the road-user, age and "
         "time-of-day subgroup analyses are exploratory: they entail multiple comparisons "
         "without formal adjustment, and the open-air road-user gradient may in part reflect "
         "weather-related differences in activity rather than physiology. Seventh, the "
         "quasi-Poisson model accounts for overdispersion but does not include an explicit "
         "autoregressive error term; although the distributed-lag, seasonal, trend and "
         "day-of-week controls absorb much short-term temporal autocorrelation, unmodelled "
         "residual serial correlation could affect confidence intervals. These findings should "
         "motivate, but cannot replace, individual-level studies linking crash decedents to "
         "ambient heat and, where available, post-mortem findings.")

    h(doc, "Conclusion", 1)
    para(doc,
         "Unusually hot days are associated with an acute, same-day excess of US traffic-crash "
         "mortality comparable to recorded direct-heat deaths, supporting the plausibility "
         "of an under-recognised heat contribution to road deaths. The findings suggest "
         "heat-aware road safety and climate adaptation planning could reduce an uncounted "
         "burden, although the Japanese data are too sparse for a firm conclusion and "
         "individual-level mechanistic evidence is still needed.")

    h(doc, "Contributors", 1)
    para(doc,
         "Tatsuki Onishi conceived the study, performed the analysis, drafted the "
         "manuscript, and approved the final version. The author had full access to all "
         "data and verified the reported results.")

    h(doc, "Declaration of competing interests", 1)
    para(doc,
         "We declare no competing interests.")

    h(doc, "Funding", 1)
    para(doc,
         "There was no specific funding for this study. The corresponding author had full "
         "access to all the data and had final responsibility for the decision to submit for "
         "publication.")

    h(doc, "Declaration of generative AI use", 1)
    para(doc,
         "During the preparation of this work the author used generative artificial "
         "intelligence tools to assist with code development and manuscript drafting. After "
         "using these tools, the author reviewed and edited the content and takes full "
         "responsibility for the published content.")

    h(doc, "Data availability", 1)
    para(doc,
         "All data are public: FARS, GHCN-Daily, CDC WONDER, US Census Bureau Population "
         "Estimates Program, FHWA Highway Statistics VM-2, FHWA/FRED, EIA and the Japanese "
         "NPA accident open data (see References for sources). Processed data files are "
         "provided in the repository so the manuscript and figures can be regenerated without "
         "API keys. The complete analysis code and reproducible pipeline (make all for raw data "
         "and figures, then make ehp for the EHP submission package) is openly available at "
         "https://github.com/bougtoir/heat-accidents-mortality; every reported number, figure "
         "and table is regenerated from source with no hard-coded values.")

    h(doc, "References", 1)
    for i, key in enumerate(_ref_order, 1):
        p = doc.add_paragraph(f"{i}. {REF_TEXT[key]}"); p.paragraph_format.space_after = Pt(4)
    missing = set(REF_TEXT) - set(_ref_order)
    assert not missing, f"uncited references: {missing}"

    path = os.path.join(MAN, filename)
    doc.save(path); print("wrote", path)


def build_pptx():
    prs = Presentation(); prs.slide_width = PInches(13.333); prs.slide_height = PInches(7.5)
    blank = prs.slide_layouts[6]
    for i, (img, cap) in enumerate(FIGURES, 1):
        s = prs.slides.add_slide(blank)
        tb = s.shapes.add_textbox(PInches(0.5), PInches(0.2), PInches(12.3), PInches(0.6))
        tf = tb.text_frame; tf.text = f"Figure {i}"
        tf.paragraphs[0].runs[0].font.size = PPt(22); tf.paragraphs[0].runs[0].font.bold = True
        s.shapes.add_picture(os.path.join(FIG, img), PInches(2.5), PInches(1.0), height=PInches(4.8))
        cb = s.shapes.add_textbox(PInches(0.5), PInches(6.0), PInches(12.3), PInches(1.3))
        ctf = cb.text_frame; ctf.word_wrap = True; ctf.text = cap
        ctf.paragraphs[0].runs[0].font.size = PPt(12)
    path = os.path.join(MAN, "figures.pptx"); prs.save(path); print("wrote", path)


def build_tables_docx():
    doc = Document(); setup(doc)
    doc.add_heading("Tables (editable)", 1)
    tbl1(doc); tbl2(doc); tbl4(doc); tbl3(doc); tbl5(doc); tbl6(doc)
    path = os.path.join(MAN, "tables.docx"); doc.save(path); print("wrote", path)


# STROBE (v4) checklist items mapped to where each is addressed in this manuscript.
STROBE_ITEMS = [
    ("1a", "Title/abstract: study design in title or abstract", "Title; Summary (Methods)"),
    ("1b", "Abstract: informative and balanced summary", "Summary (semi-structured)"),
    ("2", "Background/rationale", "Introduction"),
    ("3", "Objectives / hypotheses", "Introduction (final sentence)"),
    ("4", "Study design", "Methods (ecological time-series, quasi-Poisson DLM)"),
    ("5", "Setting, locations, periods", "Methods (US states 2016-2022; Japan prefectures 2019-2024)"),
    ("6", "Participants / units of analysis", "Methods (state-day and prefecture-day panels); Table 1"),
    ("7", "Variables (outcome, exposure, confounders)", "Methods (crash deaths; temperature anomaly; adjustments)"),
    ("8", "Data sources / measurement", "Methods; Data availability (FARS, GHCN-Daily, NPA, CDC WONDER, FHWA, EIA)"),
    ("9", "Bias", "Methods (anomaly design); Discussion (limitations)"),
    ("10", "Study size", "Table 1 (death counts, years, units)"),
    ("11", "Quantitative variables handling", "Methods (spline lag windows; anomaly construction)"),
    ("12", "Statistical methods", "Methods (quasi-Poisson DLM, attributable risk, projections, subgroups)"),
    ("13", "Descriptive data", "Results (panel sizes); Table 1"),
    ("14", "Outcome data", "Results; Tables 2-4"),
    ("15", "Main results (estimates, CIs)", "Results; Tables 2-6; Figures 1-10"),
    ("16", "Other analyses (subgroups, sensitivity)", "Results (activity, precipitation and heat-stress controls; road-user/age/time-of-day; projection); Table 6"),
    ("17", "Key results", "Discussion (first paragraph); Conclusion"),
    ("18", "Limitations", "Discussion (limitations paragraph)"),
    ("19", "Interpretation", "Discussion; Summary (Interpretation)"),
    ("20", "Generalisability", "Discussion (ecological scope; US vs Japan)"),
    ("21", "Funding", "Role of the funding source; Summary (Funding)"),
]


def build_strobe():
    doc = Document(); setup(doc)
    doc.add_heading("STROBE checklist \u2014 cross-sectional/ecological time-series study", 1)
    para(doc, "Strengthening the Reporting of Observational Studies in Epidemiology "
              "(STROBE). Item numbering follows the STROBE Statement; the third column "
              "indicates where each item is addressed in the manuscript.")
    t = doc.add_table(rows=1, cols=3); t.style = "Table Grid"
    for j, hd in enumerate(("Item", "Recommendation", "Location in manuscript")):
        run = t.rows[0].cells[j].paragraphs[0].add_run(hd); run.bold = True
    for item, rec, loc in STROBE_ITEMS:
        cells = t.add_row().cells
        cells[0].text = item; cells[1].text = rec; cells[2].text = loc
    path = os.path.join(MAN, "strobe_checklist.docx"); doc.save(path); print("wrote", path)


def build_cover_letter():
    doc = Document(); setup(doc)
    para(doc, "[PLACEHOLDER date]")
    para(doc, "The Editor-in-Chief, [Target journal]")
    para(doc, "Dear Editors,")
    para(doc,
         "We submit for your consideration our manuscript, \u201cAmbient heat as an "
         "under-recognised risk factor for US traffic-crash mortality: a distributed-lag "
         "analysis with road-safety implications\u201d, as a Research Article.")
    para(doc,
         "Using only public data, this ecological, population-level study shows that days "
         "hotter than the local seasonal norm carry an acute excess of US traffic-crash "
         "deaths that survives adjustment for national driving activity, is concentrated in "
         "heat-exposed open-air road users, and is comparable in magnitude to all officially "
         "recorded direct-heat deaths \u2014 a burden that would be invisible to routine "
         "cause-of-death surveillance and that our scenario projections suggest would grow "
         "under warming. We believe this speaks directly to the journal's focus on the "
         "environmental determinants of health and on under-recognised climate-sensitive "
         "mortality burdens.")
    para(doc,
         "The manuscript is original, is not under consideration elsewhere, and all authors "
         "approve submission. All data are public and the complete analysis pipeline is "
         "openly available and fully reproducible (make all), with no hard-coded results. We "
         "declare [PLACEHOLDER competing interests]. We suggest the manuscript has not been "
         "posted to a preprint server [PLACEHOLDER confirm].")
    para(doc, "We look forward to your assessment.")
    para(doc, "Yours sincerely,")
    para(doc, "[PLACEHOLDER corresponding author, on behalf of all authors]")
    path = os.path.join(MAN, "cover_letter.docx"); doc.save(path); print("wrote", path)


def build_submission_figures():
    """Copy each figure to a submission folder named by figure number (Lancet/EM
    uploads figures as separate files, not embedded in the manuscript)."""
    sub = os.path.join(MAN, "submission_figures")
    os.makedirs(sub, exist_ok=True)
    for i, (img, _cap) in enumerate(FIGURES, 1):
        ext = os.path.splitext(img)[1]
        shutil.copyfile(os.path.join(FIG, img), os.path.join(sub, f"Figure{i}{ext}"))
        pdf = os.path.splitext(img)[0] + ".pdf"
        pdf_src = os.path.join(FIG, pdf)
        if os.path.exists(pdf_src):
            shutil.copyfile(pdf_src, os.path.join(sub, f"Figure{i}.pdf"))
    print("wrote", sub, f"({len(FIGURES)} figures)")


def build_submission_zip():
    """Bundle the complete submission package into a single zip for delivery."""
    stage = os.path.join(MAN, "submission_package")
    if os.path.exists(stage):
        shutil.rmtree(stage)
    os.makedirs(stage)
    files = [
        "heat_crash_mortality_submission.docx",  # main text, legends only
        "tables.docx",
        "figures.pptx",
        "strobe_checklist.docx",
        "cover_letter.docx",
    ]
    for name in files:
        src = os.path.join(MAN, name)
        if os.path.exists(src):
            shutil.copyfile(src, os.path.join(stage, name))
    sub_src = os.path.join(MAN, "submission_figures")
    if os.path.isdir(sub_src):
        shutil.copytree(sub_src, os.path.join(stage, "figures"))
    zip_base = os.path.join(MAN, "submission_package")
    shutil.make_archive(zip_base, "zip", stage)
    shutil.rmtree(stage)
    print("wrote", zip_base + ".zip")


if __name__ == "__main__":
    build_manuscript()                                              # inline (reading copy)
    build_manuscript("heat_crash_mortality_submission.docx", embed=False)  # legends-only
    build_pptx(); build_tables_docx()
    build_strobe(); build_cover_letter(); build_submission_figures()
    build_submission_zip()
