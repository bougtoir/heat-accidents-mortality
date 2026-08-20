#!/usr/bin/env python3
"""Create editable PPTX with one figure per slide for Results in Engineering submission.

All 11 figures from the integrated manuscript, widescreen 13.333 x 7.5 inches.
"""

import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pathlib import Path

SCRIPT_DIR = Path(__file__).resolve().parent
# In-repo figure directories (public mirror root is dvs_noise_inverse_problem)
FIG_DIR_DVS = SCRIPT_DIR.parent          # noise_inverse_demo.py / systematic_evaluation.py outputs
FIG_DIR_SR = SCRIPT_DIR / 'sr_figures'   # generate_sr_figures.py outputs


def load_metrics():
    results_dir = SCRIPT_DIR.parent / 'results'
    with open(results_dir / 'evaluation_summary.json') as f:
        summary = json.load(f)
    with open(results_dir / 'demo_summary.json') as f:
        demo = json.load(f)
    m = summary['methods']
    return {
        'n_recordings': summary['n_recordings'],
        'fano_auc': m['fano_filter']['auc_mean'],
        'temporal_auc': m['temporal_filter']['auc_mean'],
        'fano_nrr': m['fano_filter']['nrr_mean'],
        'fano_spr': m['fano_filter']['spr_mean'],
        'demo_removal': demo['removal_pct'],
    }


FIGURES = [
    {
        'source': FIG_DIR_SR,
        'file': 'fig1_schematic.png',
        'title': 'Figure 1',
        'caption': (
            'Conceptual schematic of covariate-adjusted stochastic resonance. '
            '(a) Threshold detector: subthreshold signal + noise triggers events. '
            '(b) SR: event-rate modulation peaks at intermediate noise. '
            '(c) Covariate adjustment narrows residual, shifting the SR operating point.'
        ),
    },
    {
        'source': FIG_DIR_SR,
        'file': 'fig3_optimal_rho.png',
        'title': 'Figure 2',
        'caption': (
            'Optimal noise model accuracy. (a) \u03c1* vs. input noise: 0 in SR regime, '
            '\u221a(1 \u2212 \u03b8\u00b2/\u03c3\u00b2) in excess-noise regime. '
            '(b) SNR gain at \u03c1* grows exponentially with \u03c3/\u03b8.'
        ),
    },
    {
        'source': FIG_DIR_DVS,
        'file': 'fig2_g3_pipeline_en.png',
        'title': 'Figure 3',
        'caption': (
            'PI-DC-DVS pipeline: (1) A5 noise model + auxiliary channels; '
            '(2) Bayesian inverse solution; (3) probabilistic thinning; '
            '(4) calibration including Cal-6 satellite trails.'
        ),
    },
    {
        'source': FIG_DIR_SR,
        'file': 'fig2_sr_curves.png',
        'title': 'Figure 4',
        'caption': (
            'SR curves (A/\u03b8 = 0.3). Analytical (solid) + Monte Carlo validation '
            '(squares). Covariate adjustment shifts peak rightward.'
        ),
    },
    {
        'source': FIG_DIR_SR,
        'file': 'fig4_detection_probability.png',
        'title': 'Figure 5',
        'caption': (
            '(a) Detection probability P_D and (b) false alarm P_FA vs. noise '
            'for different \u03c1 (A/\u03b8 = 0.4). Higher \u03c1 suppresses effective noise.'
        ),
    },
    {
        'source': FIG_DIR_SR,
        'file': 'fig5_roc_comparison.png',
        'title': 'Figure 6',
        'caption': (
            'ROC curves at \u03c3_n/\u03b8 = 1.5. Covariate adjustment (\u03c1 = 0.95) '
            'achieves near-ideal signal/noise separation.'
        ),
    },
    {
        'source': FIG_DIR_DVS,
        'file': 'fig6_a5_simulation.png',
        'title': 'Figure 7',
        'caption': (
            'A5-based noise rate simulation. (a) Predicted noise rate [evt/s/pix]; '
            '(b) SNR improvement at 90% accuracy; (c) SNR vs. temperature.'
        ),
    },
    {
        'source': FIG_DIR_DVS,
        'file': 'fig5_systematic_evaluation.png',
        'title': 'Figure 8',
        'caption': (
            'Systematic evaluation on the public EBSSA dataset ({n_recordings} recordings). '
            '(a) NRR, (b) SPR, (c) F1, (d) ROC-AUC. '
            'Fano filter achieves best overall balance (AUC = {fano_auc:.3f}).'
        ),
    },
    {
        'source': FIG_DIR_SR,
        'file': 'fig7_dvs_application.png',
        'title': 'Figure 9',
        'caption': (
            'DVS results in SR framework. (a) ROC-AUC: Fano {fano_auc:.3f} vs temporal {temporal_auc:.3f}. '
            '(b) NRR vs SPR: {fano_nrr:.1%} noise removed, {fano_spr:.1%} signal preserved.'
        ),
    },
    {
        'source': FIG_DIR_DVS,
        'file': 'fig3_noise_inverse_demo.png',
        'title': 'Figure 10',
        'caption': (
            'Proof-of-concept: (a) raw events; (b) noise rate map; '
            '(c) bimodal P_noise; (d) residual after {demo_removal}% noise removal.'
        ),
    },
    {
        'source': FIG_DIR_DVS,
        'file': 'fig4_sn_improvement.png',
        'title': 'Figure 11',
        'caption': (
            'SNR improvement: (a) Fano spatial map (noise F\u22481 vs signal F\u226b1); '
            '(b) temporal dynamics; (c) per-pixel SNR distribution.'
        ),
    },
]


def _add_textbox(slide, left, top, width, height, text, font_size=Pt(12),
                 bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.bold = bold
    p.alignment = alignment
    return txBox


def build_figure_slide(prs, fig_info):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    img_path = fig_info['source'] / fig_info['file']

    # Title
    _add_textbox(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.5),
                 fig_info['title'], font_size=Pt(24), bold=True)

    if not img_path.exists():
        _add_textbox(slide, Inches(2), Inches(3), Inches(8), Inches(1),
                     f"[MISSING: {img_path}]", font_size=Pt(18))
        return slide

    from PIL import Image
    with Image.open(img_path) as img:
        img_w, img_h = img.size
    aspect = img_w / img_h

    max_w = Inches(12)
    max_h = Inches(5.0)
    if aspect > (max_w / max_h):
        w = max_w
        h = int(w / aspect)
    else:
        h = max_h
        w = int(h * aspect)

    left = int((prs.slide_width - w) / 2)
    top = Inches(0.9)
    slide.shapes.add_picture(str(img_path), left, top, w, h)

    cap_top = top + h + Inches(0.15)
    _add_textbox(slide, Inches(0.5), cap_top, Inches(12), Inches(1.2),
                 fig_info['caption'], font_size=Pt(12))

    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    metrics = load_metrics()
    for fig_info in FIGURES:
        fig_info['caption'] = fig_info['caption'].format(**metrics)
        build_figure_slide(prs, fig_info)

    out_path = SCRIPT_DIR / 'figures_rie.pptx'
    prs.save(str(out_path))
    print(f"PPTX saved: {out_path}")


if __name__ == '__main__':
    main()
